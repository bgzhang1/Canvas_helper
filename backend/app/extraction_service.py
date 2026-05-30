from __future__ import annotations

import io
import json
import zipfile
from pathlib import Path

from bs4 import BeautifulSoup

from .db import Database, utc_now


class ExtractionService:
    def __init__(
        self,
        db: Database,
        data_dir: Path,
        *,
        ocr_enabled: bool = True,
        ocr_languages: str = "eng+chi_sim",
        ocr_max_pages: int = 20,
    ):
        self.db = db
        self.data_dir = data_dir
        self.ocr_enabled = ocr_enabled
        self.ocr_languages = ocr_languages
        self.ocr_max_pages = ocr_max_pages

    async def extract_course(self, course_id: int) -> dict[str, int]:
        return await self.extract_files(course_id)

    async def extract_files(
        self,
        course_id: int,
        file_ids: list[int] | None = None,
    ) -> dict[str, int]:
        counts = {"extracted": 0, "partial": 0, "skipped": 0, "failed": 0}
        with self.db.connect() as conn:
            args: list[int] = [course_id]
            file_filter = ""
            if file_ids is not None:
                if not file_ids:
                    return counts
                placeholders = ",".join("?" for _ in file_ids)
                file_filter = f" AND id IN ({placeholders})"
                args.extend(file_ids)
            rows = conn.execute(
                f"""
                SELECT * FROM files
                WHERE course_id = ? AND local_path IS NOT NULL
                {file_filter}
                ORDER BY display_name
                """,
                args,
            ).fetchall()

        course_name = self._course_label(course_id)
        for row in rows:
            try:
                local_path = Path(row["local_path"])
                if not local_path.exists():
                    counts["skipped"] += 1
                    self.db.add_event(
                        category="file",
                        action="file_extracted",
                        status="warning",
                        title="File extraction skipped",
                        course_id=course_id,
                        course_name=course_name,
                        item_id=row["id"],
                        item_name=row["display_name"],
                        message="Cached file is missing on disk.",
                        metadata={"outcome": "skipped"},
                    )
                    continue
                text, status, warning = self.extract_file(local_path, row["content_type"])
                output = (
                    self.data_dir
                    / "extracted"
                    / f"course_{course_id}"
                    / f"{row['id']}.txt"
                )
                output.parent.mkdir(parents=True, exist_ok=True)
                output.write_text(text, encoding="utf-8", errors="replace")
                outline = self._build_outline(text, row["display_name"])
                with self.db.connect() as conn:
                    conn.execute(
                        """
                        UPDATE files
                        SET extraction_status = ?,
                            extraction_error = ?,
                            extracted_text_path = ?,
                            outline_json = ?,
                            extracted_at = ?
                        WHERE id = ?
                        """,
                        (
                            status,
                            warning,
                            str(output),
                            json.dumps(outline, ensure_ascii=False),
                            utc_now(),
                            row["id"],
                        ),
                    )
                if status == "partial":
                    counts["partial"] += 1
                elif status == "extracted":
                    counts["extracted"] += 1
                elif status == "unsupported":
                    counts["skipped"] += 1
                else:
                    counts["failed"] += 1
                self.db.add_event(
                    category="file",
                    action="file_extracted",
                    status="success" if status == "extracted" else "warning" if status in {"partial", "unsupported"} else "failed",
                    title="File extracted" if status == "extracted" else "File extraction partial" if status == "partial" else "File extraction unsupported",
                    course_id=course_id,
                    course_name=course_name,
                    item_id=row["id"],
                    item_name=row["display_name"],
                    message=warning,
                    metadata={"outcome": status},
                )
            except Exception as exc:
                with self.db.connect() as conn:
                    conn.execute(
                        """
                        UPDATE files
                        SET extraction_status = 'error', extraction_error = ?
                        WHERE id = ?
                        """,
                        (f"{exc.__class__.__name__}: {exc}", row["id"]),
                    )
                counts["failed"] += 1
                self.db.add_event(
                    category="file",
                    action="file_extracted",
                    status="failed",
                    title="File extraction failed",
                    course_id=course_id,
                    course_name=course_name,
                    item_id=row["id"],
                    item_name=row["display_name"],
                    message=f"{exc.__class__.__name__}: {exc}",
                    metadata={"outcome": "failed"},
                )
        return counts

    def _course_label(self, course_id: int) -> str:
        with self.db.connect() as conn:
            row = conn.execute(
                "SELECT name, course_code FROM courses WHERE id = ?",
                (course_id,),
            ).fetchone()
        if not row:
            return f"course_{course_id}"
        return row["course_code"] or row["name"] or f"course_{course_id}"

    def extract_file(self, path: Path, content_type: str | None) -> tuple[str, str, str | None]:
        suffix = path.suffix.lower()
        if suffix == ".pdf" or content_type == "application/pdf":
            return self._extract_pdf(path)
        if suffix == ".pptx":
            return self._extract_pptx(path)
        if suffix == ".docx":
            return self._extract_docx(path)
        if suffix in {".html", ".htm"} or content_type == "text/html":
            html = path.read_text(encoding="utf-8", errors="replace")
            return BeautifulSoup(html, "html.parser").get_text("\n"), "extracted", None
        if suffix in {".txt", ".md", ".csv", ".py", ".java", ".c", ".cpp", ".js", ".ts"}:
            return path.read_text(encoding="utf-8", errors="replace"), "extracted", None
        if suffix == ".zip":
            return self._extract_zip_listing(path), "partial", "ZIP contents listed; nested files were not extracted"
        return "", "unsupported", f"Unsupported file type: {suffix or content_type or 'unknown'}"

    def _extract_pdf(self, path: Path) -> tuple[str, str, str | None]:
        import fitz

        doc = fitz.open(path)
        chunks: list[str] = []
        ocr_warnings: list[str] = []
        for index, page in enumerate(doc):
            chunks.append(f"\n\n--- Page {index + 1} ---\n")
            chunks.append(page.get_text("text"))
            if self.ocr_enabled and index < self.ocr_max_pages:
                ocr_text, warning = self._ocr_pdf_page(page)
                if ocr_text.strip():
                    chunks.append("\n[OCR]\n" + ocr_text)
                if warning:
                    ocr_warnings.append(warning)
        warning = "; ".join(sorted(set(ocr_warnings))) or None
        return "\n".join(chunks), "partial" if warning else "extracted", warning

    def _ocr_pdf_page(self, page) -> tuple[str, str | None]:
        try:
            from PIL import Image
            import pytesseract
        except Exception as exc:
            return "", f"OCR unavailable: {exc.__class__.__name__}"
        try:
            pixmap = page.get_pixmap(matrix=page.parent.Matrix(2, 2), alpha=False)
        except AttributeError:
            import fitz

            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        image = Image.open(io.BytesIO(pixmap.tobytes("png")))
        try:
            return pytesseract.image_to_string(image, lang=self.ocr_languages), None
        except Exception as exc:
            return "", f"OCR failed: {exc.__class__.__name__}"

    def _extract_pptx(self, path: Path) -> tuple[str, str, str | None]:
        from PIL import Image
        from pptx import Presentation

        presentation = Presentation(path)
        chunks: list[str] = []
        warnings: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            chunks.append(f"\n\n--- Slide {slide_index} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
                image = getattr(shape, "image", None)
                if self.ocr_enabled and image is not None:
                    try:
                        import pytesseract

                        pil_image = Image.open(io.BytesIO(image.blob))
                        ocr_text = pytesseract.image_to_string(pil_image, lang=self.ocr_languages)
                        if ocr_text.strip():
                            chunks.append("\n[OCR]\n" + ocr_text)
                    except Exception as exc:
                        warnings.append(f"OCR failed: {exc.__class__.__name__}")
        warning = "; ".join(sorted(set(warnings))) or None
        return "\n".join(chunks), "partial" if warning else "extracted", warning

    def _extract_docx(self, path: Path) -> tuple[str, str, str | None]:
        from PIL import Image
        from docx import Document

        document = Document(path)
        chunks = [p.text for p in document.paragraphs if p.text]
        for table in document.tables:
            for row in table.rows:
                chunks.append(" | ".join(cell.text for cell in row.cells))

        warnings: list[str] = []
        if self.ocr_enabled:
            for rel in document.part._rels.values():
                if "image" not in rel.reltype:
                    continue
                try:
                    import pytesseract

                    image = Image.open(io.BytesIO(rel.target_part.blob))
                    ocr_text = pytesseract.image_to_string(image, lang=self.ocr_languages)
                    if ocr_text.strip():
                        chunks.append("\n[OCR]\n" + ocr_text)
                except Exception as exc:
                    warnings.append(f"OCR failed: {exc.__class__.__name__}")
        warning = "; ".join(sorted(set(warnings))) or None
        return "\n".join(chunks), "partial" if warning else "extracted", warning

    def _extract_zip_listing(self, path: Path) -> str:
        with zipfile.ZipFile(path) as archive:
            return "\n".join(info.filename for info in archive.infolist())

    def _build_outline(self, text: str, fallback_title: str) -> list[dict[str, str]]:
        candidates: list[str] = []
        for raw_line in text.splitlines():
            line = " ".join(raw_line.strip().split())
            if not 4 <= len(line) <= 120:
                continue
            lowered = line.lower()
            if (
                lowered.startswith(("chapter", "lecture", "week", "topic", "unit", "part "))
                or line[:1].isdigit()
                or line.isupper()
            ):
                candidates.append(line)
            if len(candidates) >= 20:
                break
        if not candidates:
            candidates = [fallback_title]
        return [{"title": item} for item in dict.fromkeys(candidates)]
